from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET = ROOT / "assets_utn_visuales"
OUT = ROOT / "informe_tecnico_utn_3d.pptx"

BG = RGBColor(7, 16, 24)
TEXT = RGBColor(226, 232, 240)
MUTED = RGBColor(148, 163, 184)
CYAN = RGBColor(56, 189, 248)
BLUE = RGBColor(14, 165, 233)
RED = RGBColor(248, 113, 113)
GREEN = RGBColor(74, 222, 128)
AMBER = RGBColor(251, 191, 36)
CARD = RGBColor(12, 28, 40)
CARD_2 = RGBColor(10, 23, 33)
BORDER = RGBColor(24, 61, 79)
WHITE = RGBColor(255, 255, 255)

W = 13.333
H = 7.5


def rgb(hex_value: str):
    hex_value = hex_value.lstrip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def shape(slide, kind, x, y, w, h, fill=None, line=None, radius=True):
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    return s


def box(slide, x, y, w, h, fill=CARD, line=BORDER):
    return shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line)


def line(slide, x1, y1, x2, y2, color=BORDER, width=1.5):
    s = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    s.line.color.rgb = color
    s.line.width = Pt(width)
    return s


def add_text(slide, x, y, w, h, value, size=18, color=TEXT, bold=False, font="Atkinson Hyperlegible", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.05):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    parts = str(value).split("\n")
    for idx, part in enumerate(parts):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = part
        p.alignment = align
        p.space_after = Pt(2)
        if p.runs:
            r = p.runs[0]
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
    return tb


def label(slide, x, y, value, color=CYAN, w=1.4):
    box(slide, x, y, w, 0.28, CARD_2, BORDER)
    add_text(slide, x + 0.08, y + 0.02, w - 0.16, 0.22, value.upper(), 8.5, color, True, "Aptos Mono", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def title(slide, kicker, heading, sub=None):
    label(slide, 0.55, 0.34, kicker, CYAN, max(1.25, 0.14 * len(kicker) + 0.32))
    add_text(slide, 0.55, 0.78, 8.9, 0.6, heading, 28, TEXT, True, "Atkinson Hyperlegible")
    if sub:
        add_text(slide, 0.58, 1.42, 9.8, 0.38, sub, 12.5, MUTED, False, "Atkinson Hyperlegible")


def footer(slide, number):
    line(slide, 0.55, 7.08, 12.78, 7.08, BORDER, 0.8)
    add_text(slide, 0.58, 7.15, 8, 0.2, "UTN FRLP · Álgebra y Geometría Analítica · Comisión S16", 7.5, MUTED, False, "Aptos Mono")
    add_text(slide, 12.15, 7.12, 0.55, 0.24, f"{number:02d}", 9, CYAN, True, "Aptos Mono", align=PP_ALIGN.RIGHT)


def bg(slide, number, motif="cyan"):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    # Quiet 3D atmosphere: horizon, micro-grid and a few distant particles.
    line(slide, 0.55, 6.72, 12.78, 6.72, RGBColor(14, 37, 50), 0.8)
    for x in [1.2, 2.55, 3.9, 5.25, 6.6, 7.95, 9.3, 10.65, 12.0]:
        line(slide, x, 6.72, x + 0.42, 6.1, RGBColor(10, 28, 39), 0.55)
    for x, y, r in [(1.1, 0.25, 0.018), (5.2, 0.24, 0.012), (8.65, 0.18, 0.015), (11.8, 0.47, 0.014), (12.5, 2.2, 0.012)]:
        s = shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x, y, r, r, CYAN if motif == "cyan" else RED, None)
    if motif == "red":
        line(slide, 10.7, 1.2, 12.8, 0.55, RGBColor(55, 29, 40), 0.8)


def picture(slide, name, x, y, w=None, h=None):
    path = ASSET / name
    if w is not None and h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), height=Inches(h))


def dot(slide, x, y, color, size=0.18):
    shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x, y, size, size, color, None)


def node(slide, x, y, n, color, text_value=None):
    shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x, y, 0.52, 0.52, color, None)
    add_text(slide, x, y + 0.02, 0.52, 0.45, n, 11, BG, True, "Aptos Mono", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if text_value:
        add_text(slide, x - 0.42, y + 0.62, 1.35, 0.34, text_value, 10, TEXT, True, "Atkinson Hyperlegible", align=PP_ALIGN.CENTER)


def make_deck():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # 1 — cover
    s = prs.slides.add_slide(blank); bg(s, 1)
    label(s, 0.58, 0.52, "INFORME TÉCNICO", CYAN, 1.7)
    add_text(s, 0.58, 1.08, 5.7, 0.62, "Álgebra y\ngeometría analítica", 30, TEXT, True)
    add_text(s, 0.58, 2.42, 5.2, 1.15, "Auditoría\nepistemológica", 44, WHITE, True, "Crimson Pro")
    add_text(s, 0.62, 3.82, 4.7, 0.5, "Recta y plano en R³ · errores humanos y desempeño de LLMs", 13.5, MUTED)
    line(s, 0.62, 4.62, 3.05, 4.62, CYAN, 2)
    add_text(s, 0.62, 4.78, 3.6, 0.42, "Santiago Sessa · Mateo Rau\nLucio Pieroni · Lucas Bazan", 10.5, TEXT)
    add_text(s, 0.62, 5.58, 3.7, 0.42, "UTN FRLP · Comisión S16\nAgosto de 2026", 9.5, MUTED, False, "Aptos Mono")
    picture(s, "08_overlay_interseccion_transparente.png", 6.55, 1.02, 6.2, 5.92)
    # Orbital ring echoing the reference page's floating 3D object.
    ring = shape(s, MSO_AUTO_SHAPE_TYPE.OVAL, 8.95, 1.45, 3.1, 3.1, None, RGBColor(18, 67, 85))
    ring.line.width = Pt(1.2)
    footer(s, 1)

    # 2 — protocol
    s = prs.slides.add_slide(blank); bg(s, 2)
    title(s, "RESUMEN EJECUTIVO", "Un protocolo en cinco fases", "La investigación combina cálculo, contraste y juicio crítico.")
    xs = [0.9, 3.25, 5.6, 7.95, 10.3]
    labs = [("01", "Resolver", "patrones\nmatemáticos", BLUE), ("02", "Contrastar", "grupo vs.\nmodelos", CYAN), ("03", "Verificar", "puntos y\nsignos", GREEN), ("04", "Tensionar", "prompts\nadversarios", RED), ("05", "Reflexionar", "aporte\npedagógico", TEXT)]
    for i in range(4): line(s, xs[i] + 0.55, 3.25, xs[i + 1] - 0.05, 3.25, RGBColor(36, 80, 96), 1.2)
    for x, (n, lab, sub, col) in zip(xs, labs):
        node(s, x, 3.0, n, col, lab)
        add_text(s, x - 0.42, 4.2, 1.35, 0.48, sub, 10, MUTED, False, "Aptos Mono", align=PP_ALIGN.CENTER)
    box(s, 3.45, 5.38, 6.45, 0.62, CARD_2, BORDER)
    add_text(s, 3.62, 5.55, 6.1, 0.25, "calcular  →  contrastar  →  verificar  →  explicar", 12, TEXT, True, "Aptos Mono", align=PP_ALIGN.CENTER)
    picture(s, "10_overlay_auditoria_nodos_transparente.png", 9.8, 0.95, 2.8, 2.15)
    footer(s, 2)

    # 3 — intersection
    s = prs.slides.add_slide(blank); bg(s, 3)
    title(s, "FASE 1 · RESOLUCIÓN", "01 / Intersección entre recta y plano", "Del planteo paramétrico al punto verificable.")
    box(s, 0.65, 2.2, 3.55, 1.05, CARD, BORDER); label(s, 0.85, 2.4, "PLANO π", BLUE, 1.0); add_text(s, 0.85, 2.77, 3.0, 0.28, "2x − y + z − 6 = 0", 14, TEXT, True, "Aptos Mono")
    box(s, 0.65, 3.48, 3.55, 1.32, CARD, BORDER); label(s, 0.85, 3.68, "RECTA r", CYAN, 1.0); add_text(s, 0.85, 4.05, 3.05, 0.58, "x = −1 + 3λ\ny = 2 + λ   z = −2λ", 13, TEXT, False, "Aptos Mono")
    box(s, 0.65, 5.22, 3.55, 0.8, RGBColor(29, 25, 32), RGBColor(98, 44, 57)); add_text(s, 0.88, 5.42, 3.0, 0.25, "I = (9, 16/3, −20/3)", 15, RED, True, "Aptos Mono"); add_text(s, 0.88, 5.72, 3.0, 0.18, "verificación cruzada ✓", 10, MUTED, False, "Aptos Mono")
    picture(s, "08_overlay_interseccion_transparente.png", 4.0, 1.95, 8.5, 4.95)
    footer(s, 3)

    # 4 — angle
    s = prs.slides.add_slide(blank); bg(s, 4)
    title(s, "FASE 1 · DEDUCCIÓN", "02 / Ángulo entre recta y plano", "La clave es distinguir el ángulo con la normal del ángulo con el plano.")
    box(s, 0.65, 2.22, 4.25, 1.55, CARD, BORDER); label(s, 0.88, 2.45, "VECTORES", CYAN, 1.1); add_text(s, 0.88, 2.84, 3.65, 0.55, "n = (1, −2, 2)\nd = (1, 2, 2)", 16, TEXT, True, "Aptos Mono")
    box(s, 0.65, 4.05, 4.25, 1.48, CARD_2, BORDER); add_text(s, 0.88, 4.28, 3.7, 0.72, "sin α = |n · d| / (‖n‖ · ‖d‖)\nα = 6,38°", 15, TEXT, True, "Aptos Mono"); add_text(s, 0.88, 5.15, 3.65, 0.2, "= 6° 22′ 46″", 11, RED, True, "Aptos Mono")
    picture(s, "09_overlay_angulo_normal_transparente.png", 5.05, 2.0, 7.9, 4.95)
    footer(s, 4)

    # 5 — parameter m
    s = prs.slides.add_slide(blank); bg(s, 5)
    title(s, "FASE 1 · DISCUSIÓN", "03 / El parámetro m no siempre existe", "Dos condiciones, dos conclusiones distintas.")
    box(s, 0.7, 2.22, 5.65, 3.65, CARD, BORDER); label(s, 0.95, 2.48, "CASO A", GREEN, 1.0); add_text(s, 0.95, 2.95, 4.9, 0.35, "d ⟂ n", 22, TEXT, True, "Crimson Pro"); add_text(s, 0.95, 3.45, 4.9, 0.68, "Producto escalar nulo\n(m)(3) + (6)(1) + (4)(−2) = 0", 14, MUTED, False, "Aptos Mono"); add_text(s, 0.95, 4.55, 4.9, 0.4, "m = 2/3", 22, GREEN, True, "Aptos Mono"); add_text(s, 0.95, 5.1, 4.9, 0.3, "solución única y bien determinada", 11, TEXT)
    box(s, 6.55, 2.22, 6.05, 3.65, RGBColor(28, 21, 28), RGBColor(88, 41, 54)); label(s, 6.8, 2.48, "CASO B", RED, 1.0); add_text(s, 6.8, 2.95, 5.3, 0.35, "d ∥ n", 22, TEXT, True, "Crimson Pro"); add_text(s, 6.8, 3.45, 5.2, 0.68, "Proporcionalidad vectorial\nm/3 = 6/1 = 4/(−2)", 14, MUTED, False, "Aptos Mono"); add_text(s, 6.8, 4.55, 5.2, 0.4, "6 ≠ −2  →  ∄ m ∈ ℝ", 22, RED, True, "Aptos Mono"); add_text(s, 6.8, 5.1, 5.2, 0.3, "el sistema es incompatible", 11, TEXT)
    footer(s, 5)

    # 6 — projecting planes
    s = prs.slides.add_slide(blank); bg(s, 6)
    title(s, "FASES 1 Y 3 · VERIFICACIÓN", "04 / Verificar antes de cerrar la ecuación", "Tres planos cartesianos. Un único error de signo detectado a tiempo.")
    cards = [("πxy", "3x + 4y − 2 = 0", "paralelo a Z", BLUE), ("πxz", "x − 4z + 18 = 0", "corregido: −18 → +18", RED), ("πyz", "y + 3z − 14 = 0", "paralelo a X", CYAN)]
    for i, (lab, eq, note, col) in enumerate(cards):
        x = 0.7 + i * 4.15
        box(s, x, 2.25, 3.75, 2.45, CARD, BORDER); label(s, x + 0.22, 2.48, lab, col, 0.8); add_text(s, x + 0.22, 2.98, 3.2, 0.34, eq, 14, TEXT, True, "Aptos Mono");
        # small 3D plane glyph
        shape(s, MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, x + 0.45, 3.55, 2.75, 0.52, RGBColor(10, 62, 82), col)
        line(s, x + 1.1, 3.38, x + 1.1, 4.42, MUTED, 1.2); line(s, x + 0.55, 4.05, x + 3.12, 4.05, MUTED, 1.2)
        add_text(s, x + 0.22, 4.35, 3.2, 0.22, note, 10, RED if col == RED else MUTED, False, "Aptos Mono")
    box(s, 2.2, 5.42, 8.95, 0.7, RGBColor(32, 23, 29), RGBColor(88, 41, 54)); add_text(s, 2.45, 5.62, 8.4, 0.25, "P(2, −1, 5)  ·  −36 ≠ 0  →  revisar término independiente", 14, RED, True, "Aptos Mono", align=PP_ALIGN.CENTER)
    footer(s, 6)

    # 7 — matrix
    s = prs.slides.add_slide(blank); bg(s, 7)
    title(s, "FASES 2 Y 3 · CONTRASTE", "Auditoría cruzada de desempeño", "La matriz compara resultado, proceso y capacidad de verificación.")
    x0, y0 = 0.7, 2.15; widths = [3.0, 4.25, 4.4]
    headers = ["EJERCICIO", "GRUPO DE ESTUDIANTES", "MODELOS DE IA"]
    for j, htxt in enumerate(headers):
        x = x0 + sum(widths[:j]); box(s, x, y0, widths[j], 0.48, RGBColor(9, 45, 61) if j else RGBColor(19, 29, 45), CYAN if j else BORDER); add_text(s, x + 0.14, y0 + 0.12, widths[j] - 0.28, 0.2, htxt, 10, TEXT, True, "Aptos Mono")
    rows = [("1 · Intersección", "OK · punto exacto", "OK · preciso", GREEN), ("2 · Ángulo", "OK · usa seno", "PARCIAL · coseno", AMBER), ("3 · Parámetro m", "OK · detecta absurdo", "ALERTA · m = 18", RED), ("4 · Proyectantes", "ALERTA · −18 → +18", "OK · omisión πyz", GREEN)]
    for i, (a, b, c, col) in enumerate(rows):
        y = y0 + 0.58 + i * 0.7
        fill = CARD if i % 2 == 0 else CARD_2
        vals = [a, b, c]
        for j, val in enumerate(vals):
            x = x0 + sum(widths[:j]); box(s, x, y, widths[j], 0.57, fill, BORDER); add_text(s, x + 0.14, y + 0.16, widths[j] - 0.28, 0.2, val, 11, TEXT if j == 0 else (col if ("ALERTA" in val or "PARCIAL" in val) else GREEN), j == 0, "Atkinson Hyperlegible" if j == 0 else "Aptos Mono")
    picture(s, "10_overlay_auditoria_nodos_transparente.png", 10.6, 5.25, 2.1, 1.4)
    footer(s, 7)

    # 8 — adversarial prompts
    s = prs.slides.add_slide(blank); bg(s, 8, "red")
    title(s, "FASE 4 · PROMPTS ADVERSARIOS", "Cuatro pruebas para forzar el error", "La resistencia aparece cuando el modelo no obedece una premisa matemática falsa.")
    center = (6.5, 4.28)
    shape(s, MSO_AUTO_SHAPE_TYPE.OVAL, 5.62, 3.4, 1.8, 1.8, CARD_2, RED); add_text(s, 5.84, 3.84, 1.36, 0.5, "NO\nCOMPLETAR", 15, RED, True, "Aptos Mono", align=PP_ALIGN.CENTER)
    prompt_cards = [(0.8, 2.3, "01", "INCOMPATIBILIDAD", "Forzar un despeje numérico de m.", "detectan 6 ≠ −2", BLUE), (8.85, 2.3, "02", "DIVISIÓN POR CERO", "Insistir con 0λ = 10.", "declaran r ∩ π = ∅", CYAN), (0.8, 5.0, "03", "FÓRMULA ERRÓNEA", "Inducir el uso directo del coseno.", "confunden normal y plano", RED), (8.85, 5.0, "04", "COMPONENTE NULA", "Dividir por dᵧ = 0.", "aíslan y = −1", TEXT)]
    for x, y, n, head, body, foot, col in prompt_cards:
        box(s, x, y, 3.65, 1.32, CARD, BORDER); label(s, x + 0.18, y + 0.16, n, col, 0.5); add_text(s, x + 0.84, y + 0.18, 2.55, 0.22, head, 9.3, col, True, "Aptos Mono"); add_text(s, x + 0.18, y + 0.57, 3.2, 0.25, body, 10.5, TEXT); add_text(s, x + 0.18, y + 0.93, 3.2, 0.2, foot, 9.5, MUTED, False, "Aptos Mono")
    footer(s, 8)

    # 9 — hallucinated evidence
    s = prs.slides.add_slide(blank); bg(s, 9, "red")
    title(s, "FASE 4 · FENÓMENO CRÍTICO DE IA", "Cuando la auditoría fabrica evidencia", "Una respuesta segura puede seguir siendo falsa si no contrasta el material fuente.")
    box(s, 0.72, 2.18, 5.35, 3.55, CARD, BORDER); label(s, 0.98, 2.45, "FABRICACIÓN", RED, 1.25); add_text(s, 0.98, 2.95, 4.65, 0.8, "La IA atribuyó a otras herramientas desarrollos matemáticos que nunca habían emitido.", 18, TEXT, True, "Crimson Pro"); add_text(s, 0.98, 4.05, 4.65, 0.72, "Ejemplo: división por cero (10/0) y un despeje ficticio atribuido a otros modelos.", 12, MUTED); box(s, 0.98, 5.0, 4.65, 0.46, RGBColor(35, 23, 30), RGBColor(88, 41, 54)); add_text(s, 1.12, 5.12, 4.35, 0.2, "La supervisión humana no es opcional.", 11.5, RED, True, "Aptos Mono", align=PP_ALIGN.CENTER)
    picture(s, "10_overlay_auditoria_nodos_transparente.png", 6.25, 2.0, 6.3, 4.0)
    add_text(s, 7.35, 5.8, 4.1, 0.25, "SYCOPHANCY  ·  ADULACIÓN ALGORÍTMICA", 9.5, RED, True, "Aptos Mono", align=PP_ALIGN.CENTER)
    footer(s, 9)

    # 10 — conclusions
    s = prs.slides.add_slide(blank); bg(s, 10)
    title(s, "FASE 5 · REFLEXIÓN METACOGNITIVA", "Tres conclusiones para el futuro ingeniero", "La herramienta acelera el cálculo; el estudiante sostiene el criterio.")
    conclusions = [("01", "Verificación geométrica", "Sustituir puntos en ecuaciones cartesianas es el filtro ineludible de validez.", CYAN), ("02", "Límites de los LLMs", "Pueden resolver y acelerar, pero fallan ante conjuntos vacíos, sesgos y evidencia inventada.", RED), ("03", "Rol del estudiante", "Pasar de ejecutor de cuentas a auditor epistemológico fortalece la competencia analítica.", GREEN)]
    for i, (n, head, body, col) in enumerate(conclusions):
        y = 2.25 + i * 1.3; line(s, 0.95, y + 0.32, 1.58, y + 0.32, col, 3); add_text(s, 0.73, y, 0.65, 0.55, n, 22, col, True, "Aptos Mono", align=PP_ALIGN.CENTER); add_text(s, 1.75, y - 0.02, 4.9, 0.3, head, 16, TEXT, True); add_text(s, 1.75, y + 0.4, 6.3, 0.5, body, 11.5, MUTED)
    picture(s, "06_ciclo_verificacion_humana.png", 8.55, 2.0, 3.65, 3.65)
    box(s, 2.3, 6.15, 8.3, 0.42, CARD_2, BORDER); add_text(s, 2.48, 6.25, 7.95, 0.2, "La competencia central no es calcular más rápido: es verificar mejor.", 12, TEXT, True, "Aptos Mono", align=PP_ALIGN.CENTER)
    footer(s, 10)

    # 11 — oral defense
    s = prs.slides.add_slide(blank); bg(s, 11)
    title(s, "DEFENSA ORAL · ESTRUCTURA GRUPAL", "El guion técnico en 11 minutos", "Cada bloque tiene una función narrativa dentro de la defensa.")
    line(s, 1.0, 3.25, 12.2, 3.25, RGBColor(35, 79, 94), 2)
    times = [(1.0, "0:00–1:45", "Santiago", "apertura + protocolo", CYAN), (3.75, "1:45–4:00", "Mateo", "intersección + ángulo", BLUE), (6.5, "4:00–6:15", "Lucio", "m + planos + signo", GREEN), (9.25, "8:00–11:30", "Lucas", "prompts + cierre", RED)]
    for x, tm, person, desc, col in times:
        dot(s, x, 3.05, col, 0.4); add_text(s, x - 0.35, 2.42, 1.2, 0.22, tm, 9.3, col, True, "Aptos Mono", align=PP_ALIGN.CENTER); add_text(s, x - 0.45, 3.75, 1.35, 0.25, person, 13, TEXT, True, "Crimson Pro", align=PP_ALIGN.CENTER); add_text(s, x - 0.62, 4.1, 1.7, 0.38, desc, 9.5, MUTED, False, "Aptos Mono", align=PP_ALIGN.CENTER)
    box(s, 2.25, 5.35, 8.85, 0.58, CARD, BORDER); add_text(s, 2.5, 5.52, 8.35, 0.2, "problema   →   resolución   →   auditoría   →   reflexión", 13, TEXT, True, "Aptos Mono", align=PP_ALIGN.CENTER)
    footer(s, 11)

    # 12 — close
    s = prs.slides.add_slide(blank); bg(s, 12)
    title(s, "DEFENSA ANTE EL TRIBUNAL DOCENTE", "Banco de preguntas · cierre", "Dos preguntas que condensan la defensa.")
    questions = [("¿Por qué seno en lugar de coseno?", "El producto escalar mide el ángulo β con la normal. Como α + β = 90°, cos β = sin α."), ("¿Por qué alcanza una desigualdad?", "La colinealidad exige una única constante k. Si kᵧ = 6 y k𝓏 = −2, el sistema es incompatible.")]
    for i, (q, a) in enumerate(questions):
        y = 2.25 + i * 1.65; box(s, 0.7, y, 7.2, 1.25, CARD, BORDER); add_text(s, 0.95, y + 0.22, 6.7, 0.26, q, 15, CYAN, True, "Crimson Pro"); add_text(s, 0.95, y + 0.63, 6.55, 0.42, a, 10.8, MUTED)
    picture(s, "08_overlay_interseccion_transparente.png", 8.05, 2.1, 4.55, 3.1)
    add_text(s, 8.3, 5.35, 4.0, 0.45, "¡Muchas gracias\npor su atención!", 25, WHITE, True, "Crimson Pro", align=PP_ALIGN.CENTER)
    add_text(s, 8.3, 6.03, 4.0, 0.25, "Quedamos a disposición para preguntas.", 10.5, MUTED, False, "Aptos Mono", align=PP_ALIGN.CENTER)
    footer(s, 12)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    make_deck()
